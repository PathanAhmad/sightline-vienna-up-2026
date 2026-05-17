"""Build a clean, professional pitch deck.

Design language: large/small type contrast, generous whitespace, single warm
accent, asymmetric left-aligned layouts, big-number hero callouts, one idea per
slide. 8 slides, 16:9. Run from repo root:

    python scripts/build_pitch_deck.py

Output: pitch_deck.pptx in repo root.
"""
from __future__ import annotations

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.oxml.ns import qn
from pptx.util import Inches, Pt, Emu

# ---------- design tokens ----------
SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)

# palette
BG = RGBColor(0xFF, 0xFF, 0xFF)
INK = RGBColor(0x0A, 0x0A, 0x0A)         # near-black headlines
BODY = RGBColor(0x33, 0x33, 0x33)        # body
MUTED = RGBColor(0x8A, 0x8A, 0x8A)       # captions / subtitles
RULE = RGBColor(0xE8, 0xE8, 0xE8)        # hairlines
ACCENT = RGBColor(0xE2, 0x58, 0x22)      # warning-tape orange, deepened
ACCENT_SOFT = RGBColor(0xFB, 0xE9, 0xDF)

# RAG palette (used once, in the solution slide)
RAG_GREEN = RGBColor(0x2A, 0x8A, 0x4F)
RAG_AMBER = RGBColor(0xD9, 0xA0, 0x2C)
RAG_RED = RGBColor(0xC8, 0x36, 0x2F)

FONT = "Calibri"  # universally installed on Windows; clean modern sans

# grid
LEFT = Inches(0.7)
RIGHT_LIMIT = Inches(12.7)
EYEBROW_Y = Inches(0.55)
HEADLINE_Y = Inches(1.35)
SUB_Y = Inches(2.85)
CONTENT_Y = Inches(4.1)
FOOTER_Y = Inches(7.05)


# ---------- helpers ----------

def add_blank(prs: Presentation):
    layout = prs.slide_layouts[6]
    s = prs.slides.add_slide(layout)
    bg = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SLIDE_W, SLIDE_H)
    bg.line.fill.background()
    bg.fill.solid()
    bg.fill.fore_color.rgb = BG
    bg.shadow.inherit = False
    return s


def _apply_run_font(run, *, size, bold=False, color=BODY, tracking=0, italic=False):
    f = run.font
    f.name = FONT
    f.size = Pt(size)
    f.bold = bold
    f.italic = italic
    f.color.rgb = color
    if tracking:
        rPr = run._r.get_or_add_rPr()
        rPr.set("spc", str(tracking))


def add_text(slide, text, left, top, width, height, *, size=22, bold=False,
             color=BODY, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP, spacing=1.15,
             tracking=0, italic=False):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    tf.vertical_anchor = anchor
    p = tf.paragraphs[0]
    p.alignment = align
    p.line_spacing = spacing
    r = p.add_run()
    r.text = text
    _apply_run_font(r, size=size, bold=bold, color=color, tracking=tracking,
                    italic=italic)
    return box


def add_bullets(slide, lines, left, top, width, height, *, size=20, color=BODY,
                spacing=1.45):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = PP_ALIGN.LEFT
        p.line_spacing = spacing
        p.space_after = Pt(8)
        r_dot = p.add_run()
        r_dot.text = "—  "
        _apply_run_font(r_dot, size=size, bold=True, color=ACCENT)
        r = p.add_run()
        r.text = line
        _apply_run_font(r, size=size, color=color)
    return box


def add_rect(slide, left, top, width, height, color):
    rect = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    rect.line.fill.background()
    rect.fill.solid()
    rect.fill.fore_color.rgb = color
    return rect


def add_rule(slide, left, top, width, *, height=Emu(12700), color=RULE):
    return add_rect(slide, left, top, width, height, color)


def add_eyebrow(slide, text, *, y=EYEBROW_Y):
    add_rect(slide, LEFT, y + Pt(8), Inches(0.32), Pt(3), ACCENT)
    add_text(slide, text, LEFT + Inches(0.5), y, Inches(8), Inches(0.35),
             size=10, bold=True, color=ACCENT, tracking=500)


def footer(slide, page_no, total):
    add_rule(slide, LEFT, FOOTER_Y - Pt(8), Inches(12.0), color=RULE)
    add_text(slide, "Sightline  ·  Photo Compliance Audit  ·  Vienna UP 2026",
             LEFT, FOOTER_Y, Inches(9), Inches(0.3),
             size=9, color=MUTED, tracking=200)
    add_text(slide, f"{page_no:02d} / {total:02d}",
             Inches(11.7), FOOTER_Y, Inches(1.0), Inches(0.3),
             size=9, color=MUTED, align=PP_ALIGN.RIGHT, tracking=200)


# ---------- slides ----------

def slide_title(prs, total):
    s = add_blank(prs)
    # left-edge accent stripe (full-height)
    add_rect(s, 0, 0, Inches(0.18), SLIDE_H, ACCENT)

    add_text(s, "VIENNA UP  ·  EUROPE TECH HACKATHON 2026",
             Inches(0.9), Inches(0.7), Inches(11), Inches(0.4),
             size=11, bold=True, color=MUTED, tracking=600)

    add_text(s, "424,000 photos.",
             Inches(0.9), Inches(2.2), Inches(12), Inches(1.3),
             size=72, bold=True, color=INK, tracking=-20)
    add_text(s, "Nobody has time to look.",
             Inches(0.9), Inches(3.35), Inches(12), Inches(1.3),
             size=72, bold=True, color=ACCENT, tracking=-20)

    add_rule(s, Inches(0.9), Inches(5.2), Inches(0.6), color=ACCENT)
    add_text(s, "Sightline  —  Photo Compliance Audit for Trench Operators",
             Inches(0.9), Inches(5.4), Inches(12), Inches(0.5),
             size=22, color=BODY)
    add_text(s, "Challenge 2  ·  AI-powered evidence review for buried-cable projects",
             Inches(0.9), Inches(5.95), Inches(12), Inches(0.4),
             size=14, color=MUTED, tracking=150)

    add_text(s, "Team  ·  [your names here]",
             Inches(0.9), Inches(6.85), Inches(11), Inches(0.4),
             size=10, color=MUTED, tracking=300)
    return s


def slide_problem(prs, total, page):
    s = add_blank(prs)
    add_eyebrow(s, "01  ·  THE PROBLEM")
    add_text(s, "Three to five days.",
             LEFT, HEADLINE_Y, Inches(12), Inches(1.2),
             size=60, bold=True, color=INK, tracking=-15)
    add_text(s, "Per project section. By hand. One operator has hundreds.",
             LEFT, Inches(2.55), Inches(12), Inches(0.7),
             size=24, color=MUTED, spacing=1.3)

    add_bullets(s, [
        "Manual photo review is the bottleneck — so it gets skipped",
        "Undocumented cables surface years later, when someone digs",
        "Same photos get reused across jobs. Nobody notices.",
    ], LEFT, Inches(4.15), Inches(11.5), Inches(2.3), size=22)

    footer(s, page, total)
    return s


def slide_solution(prs, total, page):
    s = add_blank(prs)
    add_eyebrow(s, "02  ·  WHAT WE BUILT")
    add_text(s, "A colored map of every trench.",
             LEFT, HEADLINE_Y, Inches(12), Inches(1.2),
             size=54, bold=True, color=INK, tracking=-15)
    add_text(s, "Click a stretch — see the photos, the gaps, and the reason for the color.",
             LEFT, Inches(2.7), Inches(12), Inches(0.7),
             size=20, color=MUTED, spacing=1.3)

    # trench bar visual — colored segments along a "trench line"
    bar_y = Inches(4.05)
    bar_h = Inches(0.55)
    bar_x = LEFT
    total_w = Inches(12.0)
    # six segments with mixed colors
    spans = [
        (RAG_GREEN, 0.22),
        (RAG_GREEN, 0.18),
        (RAG_AMBER, 0.13),
        (RAG_GREEN, 0.17),
        (RAG_RED, 0.10),
        (RAG_GREEN, 0.20),
    ]
    cursor = bar_x
    for color, frac in spans:
        seg_w = Emu(int(total_w * frac))
        add_rect(s, cursor, bar_y, seg_w, bar_h, color)
        cursor = Emu(cursor + seg_w)
    # tick marks at start and end
    add_text(s, "meter 0", bar_x, bar_y + bar_h + Inches(0.1),
             Inches(2), Inches(0.3), size=10, color=MUTED, tracking=200)
    add_text(s, "meter 280",
             bar_x + total_w - Inches(2), bar_y + bar_h + Inches(0.1),
             Inches(2), Inches(0.3), size=10, color=MUTED, tracking=200,
             align=PP_ALIGN.RIGHT)

    # legend
    legend_y = Inches(5.35)
    for i, (label, color) in enumerate([
        ("GREEN  ·  documented, all checks pass", RAG_GREEN),
        ("YELLOW  ·  gap > 5m, or a check failed", RAG_AMBER),
        ("RED  ·  almost no photos — look here first", RAG_RED),
    ]):
        x = LEFT
        y = legend_y + Inches(i * 0.42)
        chip = s.shapes.add_shape(MSO_SHAPE.OVAL, x, y + Inches(0.07),
                                   Inches(0.22), Inches(0.22))
        chip.line.fill.background()
        chip.fill.solid()
        chip.fill.fore_color.rgb = color
        add_text(s, label, x + Inches(0.42), y, Inches(11), Inches(0.4),
                 size=16, color=BODY)

    footer(s, page, total)
    return s


def slide_demo(prs, total, page):
    s = add_blank(prs)
    add_eyebrow(s, "03  ·  LIVE")
    add_text(s, "Demo.",
             LEFT, Inches(2.3), Inches(12), Inches(3.0),
             size=220, bold=True, color=INK, tracking=-40,
             anchor=MSO_ANCHOR.MIDDLE)
    add_rule(s, LEFT, Inches(5.7), Inches(1.2), color=ACCENT)
    add_text(s, "switch tabs  ·  talk over the click-through  ·  60 seconds",
             LEFT, Inches(5.9), Inches(12), Inches(0.4),
             size=12, color=MUTED, tracking=400)
    footer(s, page, total)
    return s


def slide_features(prs, total, page):
    s = add_blank(prs)
    add_eyebrow(s, "04  ·  WHAT THE TOOL CATCHES")
    add_text(s, "Eight checks per photo.",
             LEFT, HEADLINE_Y, Inches(12), Inches(1.2),
             size=54, bold=True, color=INK, tracking=-15)
    add_text(s, "Six read by Claude vision. Two read by math. Two filters decide which apply.",
             LEFT, Inches(2.7), Inches(12), Inches(0.7),
             size=20, color=MUTED, spacing=1.3)

    col_y = Inches(4.05)
    col_w = Inches(5.7)

    # column 1: vision
    add_text(s, "READ BY CLAUDE  ·  6", LEFT, col_y, col_w, Inches(0.4),
             size=11, bold=True, color=ACCENT, tracking=500)
    add_bullets(s, [
        "Orange warning tape visible",
        "Sand under the cable",
        "Side-view of the trench",
        "Depth reference (ruler) in frame",
        "Cable ends sealed (white caps)",
        "Faces / plates  —  privacy flag",
    ], LEFT, col_y + Inches(0.5), col_w, Inches(3.0), size=17)

    # column 2: math + filters
    col2_x = Inches(7.0)
    add_text(s, "READ BY MATH  ·  2", col2_x, col_y, col_w, Inches(0.4),
             size=11, bold=True, color=ACCENT, tracking=500)
    add_bullets(s, [
        "Duplicate fingerprint  —  catches reused photos across jobs",
        "GPS vs printed address  —  catches wrong-location submissions",
    ], col2_x, col_y + Inches(0.5), col_w, Inches(1.4), size=17)

    add_text(s, "FILTERS  ·  2", col2_x, col_y + Inches(1.85), col_w, Inches(0.4),
             size=11, bold=True, color=ACCENT, tracking=500)
    add_bullets(s, [
        "Is this photo relevant at all?",
        "What stage of work — only relevant checks fire",
    ], col2_x, col_y + Inches(2.35), col_w, Inches(1.2), size=17)

    footer(s, page, total)
    return s


def slide_numbers(prs, total, page):
    s = add_blank(prs)
    add_eyebrow(s, "05  ·  THE PILOT  ·  AND HOW WE KNOW")
    add_text(s, "3,929 photos. $15. 30 minutes.",
             LEFT, HEADLINE_Y, Inches(12), Inches(1.2),
             size=52, bold=True, color=INK, tracking=-15)
    add_text(s,
             "End-to-end on a laptop. Then measured against 214 hand-labeled photos, blind.",
             LEFT, Inches(2.7), Inches(12), Inches(0.7),
             size=20, color=MUTED, spacing=1.3)

    # three big-number callouts in a grid
    grid_y = Inches(4.1)
    col_w = Inches(4.0)
    gap = Inches(0.13)
    callouts = [
        ("~600", "duplicates caught", "before any AI cost  ·  photo fingerprints"),
        ("2,983", "stretches scored", "GREEN / YELLOW / RED  ·  every 5 m"),
        ("96%", "on safety-critical shots",
         "depth-evidence agreement  ·  72% across both phase classes  ·  n=214"),
    ]
    for i, (big, mid, sub) in enumerate(callouts):
        x = LEFT + Emu(int(i * (col_w + gap)))
        # subtle background card
        add_rect(s, x, grid_y, col_w, Inches(2.3), RGBColor(0xFA, 0xFA, 0xFA))
        # accent stripe top
        add_rect(s, x, grid_y, Inches(0.4), Pt(3), ACCENT)
        add_text(s, big, x + Inches(0.3), grid_y + Inches(0.25),
                 col_w - Inches(0.6), Inches(1.1),
                 size=64, bold=True, color=INK, tracking=-25)
        add_text(s, mid, x + Inches(0.3), grid_y + Inches(1.35),
                 col_w - Inches(0.6), Inches(0.4),
                 size=18, bold=True, color=BODY)
        add_text(s, sub, x + Inches(0.3), grid_y + Inches(1.75),
                 col_w - Inches(0.6), Inches(0.6),
                 size=11, color=MUTED, tracking=150, spacing=1.25)

    footer(s, page, total)
    return s


def slide_scale(prs, total, page):
    s = add_blank(prs)
    add_eyebrow(s, "06  ·  THE NEXT MILE")
    add_text(s, "Full operator backlog: ~$1,900.",
             LEFT, HEADLINE_Y, Inches(12), Inches(1.2),
             size=54, bold=True, color=INK, tracking=-15)
    add_text(s, "Same code. Days of laptop time. No rewrite.",
             LEFT, Inches(2.7), Inches(12), Inches(0.7),
             size=22, color=MUTED, spacing=1.3)

    add_bullets(s, [
        "Killer move  —  on-site upload catches gaps while the trench is still open",
        "Plays nicely with RTK survey verification as a phase 2 add-on",
        "What we deliberately don't pretend  —  depth-in-cm, face-blurring, custom training",
    ], LEFT, Inches(4.1), Inches(12.0), Inches(2.4), size=21)

    footer(s, page, total)
    return s


def slide_closer(prs, total, page):
    s = add_blank(prs)
    add_rect(s, 0, 0, Inches(0.18), SLIDE_H, ACCENT)
    add_text(s, "Thank you.",
             Inches(0.9), Inches(2.9), Inches(12), Inches(1.5),
             size=86, bold=True, color=INK, tracking=-20,
             anchor=MSO_ANCHOR.MIDDLE)
    add_rule(s, Inches(0.9), Inches(4.5), Inches(0.6), color=ACCENT)
    add_text(s, "Questions  —  and the colored map is still on screen.",
             Inches(0.9), Inches(4.7), Inches(12), Inches(0.5),
             size=18, color=BODY)
    add_text(s, "Vienna UP  ·  Challenge 2  ·  Sightline  —  Photo Compliance Audit",
             Inches(0.9), Inches(5.4), Inches(12), Inches(0.4),
             size=11, color=MUTED, tracking=400)
    return s


# ---------- build ----------

def build(out_path: Path):
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H

    total = 8
    slide_title(prs, total)
    slide_problem(prs, total, 2)
    slide_solution(prs, total, 3)
    slide_demo(prs, total, 4)
    slide_features(prs, total, 5)
    slide_numbers(prs, total, 6)
    slide_scale(prs, total, 7)
    slide_closer(prs, total, 8)

    prs.save(out_path)
    print(f"[deck] wrote {out_path}")


if __name__ == "__main__":
    out = Path(__file__).resolve().parent.parent / "pitch_deck.pptx"
    build(out)
