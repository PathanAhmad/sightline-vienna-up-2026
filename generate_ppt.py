"""Generate Sightline_Pitch.pptx — run with: uv run python generate_ppt.py

4-slide editorial deck for the Sunday Vienna UP pitch (3 minutes total).
Demo eats ~90 s and lives between slide 2 and slide 3 (alt-tab to the
dashboard, alt-tab back).

Arc:
  1. HOOK       (quote)               — the line they remember
  2. PIPELINE   (6 stages)            — technical credibility
  → live demo on the dashboard
  3. NUMBERS    (dataset + dial)      — two tiers, speaker carries the rest
  4. CLOSE      (echo)                — the invitation

Slides scaffold the speech — they don't replicate it. The speaker carries
depth, tone, and examples. Palette + font mirror src/ui/tokens.py so the
slide-to-demo handoff feels like one product.
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE


# ── Pilot numbers (projected from 214-photo benchmark, scaled to ~3,400
# unique photos after dedup of the 3,929-photo Maria Rain dataset) ──────────
DATASET_SIZE = "3,929 photos"
FAST_COST    = "$12"
FAST_TIME    = "90 min"
FAST_ACC     = "87% depth-evidence"
CAREFUL_COST = "$40"
CAREFUL_TIME = "3 h"
CAREFUL_ACC  = "96% depth-evidence"


# ── Palette (mirrors src/ui/tokens.py) ──────────────────────────────────────
BG       = RGBColor(0xF6, 0xF7, 0xF9)   # --c-bg
BORDER   = RGBColor(0xE5, 0xE7, 0xEB)   # --c-border
INK      = RGBColor(0x0F, 0x17, 0x2A)   # --c-text          slate-900
BODY     = RGBColor(0x47, 0x55, 0x69)   # --c-text-2        slate-600
MUTED    = RGBColor(0x64, 0x74, 0x8B)   # --c-muted         slate-500
ACCENT   = RGBColor(0x03, 0x69, 0xA1)   # --c-accent        sky-700


# ── Typography (strict 5-tier scale, do not add a sixth) ────────────────────
FONT     = "Segoe UI"
T_HERO   = 48   # the quote, the close echo, the $15 anchor
T_HEAD   = 32   # standard slide headline (pipeline stage labels)
T_BODY   = 22   # body-emphasis: bullets, supporting lines
T_CAP    = 14   # captions
T_META   = 11   # eyebrow, footer, brand mark, tech footnote


# 16:9
SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)

# Editorial margins
LEFT_X    = Inches(0.95)
CONTENT_W = Inches(11.4)


# ── Primitives ──────────────────────────────────────────────────────────────
def _prs() -> Presentation:
    p = Presentation()
    p.slide_width  = SLIDE_W
    p.slide_height = SLIDE_H
    return p


def _blank(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    bg = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SLIDE_W, SLIDE_H)
    bg.fill.solid()
    bg.fill.fore_color.rgb = BG
    bg.line.fill.background()
    bg.shadow.inherit = False
    # Thin left brand stripe -- the deck's only persistent visual mark.
    stripe = s.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, 0, 0, Inches(0.10), SLIDE_H
    )
    stripe.fill.solid()
    stripe.fill.fore_color.rgb = ACCENT
    stripe.line.fill.background()
    return s


def _text(
    slide, txt, x, y, w, h,
    size=T_BODY, color=BODY, bold=False, italic=False,
    align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP,
    tracking=0, leading=1.2,
):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = 0
    tf.margin_top = tf.margin_bottom = 0
    tf.vertical_anchor = anchor

    lines = txt.split("\n") if isinstance(txt, str) else txt
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        p.line_spacing = leading
        r = p.add_run()
        r.text = line
        f = r.font
        f.name = FONT
        f.size = Pt(size)
        f.bold = bold
        f.italic = italic
        f.color.rgb = color
        if tracking:
            rPr = r._r.get_or_add_rPr()
            rPr.set("spc", str(tracking))
    return tb


def _hairline(slide, x, y, w, color=BORDER, weight=0.5):
    line = slide.shapes.add_connector(1, x, y, x + w, y)
    line.line.color.rgb = color
    line.line.width = Pt(weight)
    return line


def _eyebrow(slide, txt):
    _text(
        slide, txt,
        LEFT_X, Inches(0.75), CONTENT_W, Inches(0.35),
        size=T_META, color=ACCENT, bold=True, tracking=400,
    )


def _footer(slide, text="SIGHTLINE  ·  VIENNA UP 2026"):
    _text(
        slide, text,
        LEFT_X, Inches(7.05), CONTENT_W, Inches(0.3),
        size=T_META, color=MUTED, tracking=300,
    )


# ── Slide 1 — the hook ──────────────────────────────────────────────────────
def slide_hook(prs):
    s = _blank(prs)

    _text(
        s, "SIGHTLINE",
        LEFT_X, Inches(0.75), CONTENT_W, Inches(0.35),
        size=T_META, color=INK, bold=True, tracking=600,
    )

    _text(
        s, "“The trench gets filled in.",
        LEFT_X, Inches(2.65), CONTENT_W, Inches(1.1),
        size=T_HERO, color=INK, bold=True, leading=1.1,
    )
    _text(
        s, "The questions never do.”",
        LEFT_X, Inches(3.85), CONTENT_W, Inches(1.1),
        size=T_HERO, color=ACCENT, bold=True, leading=1.1,
    )

    _footer(s, "VIENNA UP 2026  ·  CHALLENGE 2")
    return s


# ── Slide 2 — the pipeline ──────────────────────────────────────────────────
def slide_pipeline(prs):
    """Six-stage pipeline shown as a horizontal flow.

    Speaker elaborates on the two surprises (dedup-before-AI, GPS-from-
    printed-overlay) and lets the rest read.
    """
    s = _blank(prs)
    _eyebrow(s, "01  —  THE PIPELINE")

    _text(
        s, "Six stages. Two surprised us.",
        LEFT_X, Inches(1.4), CONTENT_W, Inches(0.9),
        size=T_HEAD, color=INK, bold=True, leading=1.15,
    )

    # Six pipeline nodes, evenly distributed across the content width.
    stages = [
        ("01", "INGEST",  "photos +\nGeoJSON route"),
        ("02", "DEDUP",   "perceptual hash\non laptop, free"),
        ("03", "READ",    "vision AI reads\nthe printed overlay"),
        ("04", "SCORE",   "7 compliance checks\nper photo, 1 call"),
        ("05", "MAP",     "every photo →\n5 m trench segment"),
        ("06", "REPORT",  "CSV + PDF\ndeficiency log"),
    ]

    n = len(stages)
    # Column layout
    col_w = CONTENT_W / n
    top_y = Inches(3.2)
    label_h = Inches(0.45)
    name_h  = Inches(0.55)
    body_h  = Inches(1.0)

    for i, (num, name, body) in enumerate(stages):
        x = LEFT_X + col_w * i
        # number badge
        _text(
            s, num,
            x, top_y, col_w, label_h,
            size=T_META, color=ACCENT, bold=True, tracking=300,
            align=PP_ALIGN.LEFT,
        )
        # stage name
        _text(
            s, name,
            x, top_y + label_h, col_w, name_h,
            size=T_BODY, color=INK, bold=True, tracking=200,
        )
        # short subtitle
        _text(
            s, body,
            x, top_y + label_h + name_h + Inches(0.05),
            col_w, body_h,
            size=T_CAP, color=BODY, leading=1.35,
        )
        # arrow to next stage (skip last)
        if i < n - 1:
            arrow_y = top_y + label_h + Inches(0.18)
            _text(
                s, "→",
                x + col_w - Inches(0.12), arrow_y,
                Inches(0.2), Inches(0.3),
                size=T_BODY, color=ACCENT, bold=True,
            )

    # Demo handoff line at the bottom.
    _hairline(s, LEFT_X, Inches(6.05), Inches(11.4))
    _text(
        s, "Live demo, next.",
        LEFT_X, Inches(6.20), CONTENT_W, Inches(0.5),
        size=T_CAP, color=ACCENT, italic=True, bold=True,
    )

    _footer(s, "02 / 04")
    return s


# ── Slide 3 — the numbers (dataset + dial) ──────────────────────────────────
def slide_numbers(prs):
    """Dataset size on top, two-tier dial underneath. The speaker delivers
    the projection caveat and the accuracy beat verbally; the slide just
    anchors the two-tier shape so judges can compare with their eyes."""
    s = _blank(prs)
    _eyebrow(s, "02  —  THE NUMBERS")

    # Dataset headline — sets the scale of what the dial operates on.
    _text(
        s, DATASET_SIZE,
        LEFT_X, Inches(1.55), CONTENT_W, Inches(1.1),
        size=T_HERO, color=INK, bold=True,
        align=PP_ALIGN.CENTER, leading=1.0,
    )
    _text(
        s, "Two ways to run it.",
        LEFT_X, Inches(2.55), CONTENT_W, Inches(0.5),
        size=T_BODY, color=MUTED, italic=True,
        align=PP_ALIGN.CENTER,
    )

    # Two-column dial — left = fast, right = careful. Hairline divider
    # between them so the contrast reads at-a-glance.
    col_w = Inches(5.4)
    gap   = Inches(0.6)
    total_w = col_w * 2 + gap
    col_x_l = LEFT_X + (CONTENT_W - total_w) / 2
    col_x_r = col_x_l + col_w + gap
    col_y   = Inches(3.55)

    # Vertical hairline between the two columns.
    _hairline(
        s,
        col_x_l + col_w + gap / 2,
        col_y + Inches(0.15),
        Inches(0),
        color=BORDER, weight=0.5,
    )
    # Manual vertical line (connector with same x for start/end is a point;
    # use a thin rectangle instead so it actually shows on PowerPoint).
    div = s.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        col_x_l + col_w + gap / 2 - Inches(0.005),
        col_y + Inches(0.20),
        Inches(0.01),
        Inches(2.2),
    )
    div.fill.solid()
    div.fill.fore_color.rgb = BORDER
    div.line.fill.background()

    def _tier_column(x, tier_label, cost, time, accuracy):
        # Eyebrow label
        _text(
            s, tier_label,
            x, col_y, col_w, Inches(0.4),
            size=T_META, color=ACCENT, bold=True, tracking=400,
            align=PP_ALIGN.CENTER,
        )
        # Cost as the hero number
        _text(
            s, cost,
            x, col_y + Inches(0.45), col_w, Inches(1.4),
            size=84, color=INK, bold=True,
            align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE,
            leading=1.0,
        )
        # Time underneath, sub-dominant
        _text(
            s, time,
            x, col_y + Inches(1.85), col_w, Inches(0.5),
            size=T_HEAD, color=BODY, bold=False,
            align=PP_ALIGN.CENTER,
        )
        # Accuracy callout (smallest, the credibility line the speaker
        # elaborates on).
        _text(
            s, accuracy,
            x, col_y + Inches(2.45), col_w, Inches(0.4),
            size=T_CAP, color=MUTED, italic=True,
            align=PP_ALIGN.CENTER,
        )

    _tier_column(col_x_l, "FAST TIER",    FAST_COST,    FAST_TIME,    FAST_ACC)
    _tier_column(col_x_r, "CAREFUL TIER", CAREFUL_COST, CAREFUL_TIME, CAREFUL_ACC)

    # Manual-baseline anchor — the business-model angle, quiet. Speaker
    # turns this into the "weeks of senior review → price of a lunch"
    # contrast; the slide just lets judges multiply in their heads.
    _text(
        s, "Manual review  ·  3 – 5 engineer-days per section",
        LEFT_X, Inches(6.30), CONTENT_W, Inches(0.35),
        size=T_CAP, color=MUTED, italic=True,
        align=PP_ALIGN.CENTER, tracking=200,
    )

    # Bottom rail — the product point.
    _hairline(s, LEFT_X + Inches(5.2), Inches(6.75), Inches(3.0),
              color=ACCENT, weight=1.5)
    _text(
        s, "Same pipeline. Operator picks.",
        LEFT_X, Inches(6.85), CONTENT_W, Inches(0.4),
        size=T_BODY, color=BODY, italic=True,
        align=PP_ALIGN.CENTER,
    )

    _footer(s, "03 / 04")
    return s


# ── Slide 4 — close (echoes the hook) ───────────────────────────────────────
def slide_close(prs):
    s = _blank(prs)

    # No eyebrow on the close -- this is the final beat, not a step.
    _text(
        s, "Don’t let the questions",
        LEFT_X, Inches(2.55), CONTENT_W, Inches(1.1),
        size=T_HERO, color=INK, bold=True, leading=1.1,
    )
    _text(
        s, "outlive the trench.",
        LEFT_X, Inches(3.75), CONTENT_W, Inches(1.1),
        size=T_HERO, color=ACCENT, bold=True, leading=1.1,
    )

    _hairline(s, LEFT_X, Inches(5.40), Inches(0.5), color=INK, weight=1.0)
    _text(
        s,
        "Sightline.  The demo’s live.",
        LEFT_X, Inches(5.60), CONTENT_W, Inches(0.5),
        size=T_BODY, color=INK, bold=True,
    )

    _footer(s, "TEAM SIGHTLINE  ·  VIENNA UP 2026")
    return s


# ── Build ───────────────────────────────────────────────────────────────────
if __name__ == "__main__":
    deck = _prs()
    slide_hook(deck)
    slide_pipeline(deck)
    slide_numbers(deck)
    slide_close(deck)

    out = "Sightline_Pitch.pptx"
    deck.save(out)
    print(f"Saved -> {out}  ({len(deck.slides)} slides)")
